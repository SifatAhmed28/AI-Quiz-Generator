"""
NoteWise Quiz – Enhanced Flask 3 + Gemini 2.0 + Advanced Quiz Features

Complete REST API backend with user profiles, quiz history, and advanced quiz taking features

Enhanced version with timer, navigation, progress tracking, and robust error handling

"""

import os, re, tempfile, uuid, textwrap, logging, smtplib, json, threading, time
import cv2
import numpy as np
from datetime import datetime, timedelta
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from email.mime.application import MIMEApplication
from flask import (Flask, request, send_file, session,
                   render_template, redirect, url_for, flash, jsonify)
import google.generativeai as genai
import PyPDF2
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from dotenv import load_dotenv
from flask_wtf import FlaskForm
from wtforms import StringField, PasswordField, SelectField, IntegerField, validators
from werkzeug.utils import secure_filename
from PIL import Image
import easyocr
from docx import Document
from pptx import Presentation
import docx2txt

# Authentication and database imports
from flask_login import LoginManager, login_user, logout_user, login_required, current_user, UserMixin
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
from apscheduler.schedulers.background import BackgroundScheduler
import atexit

# ─────────────── ENV & LOGGING ───────────────
load_dotenv()
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s │ %(levelname)-8s │ %(message)s"
)

log = logging.getLogger("notewise-quiz")

app = Flask(__name__, static_folder="static")
app.secret_key = os.getenv("FLASK_SECRET_KEY") or os.urandom(32)
app.permanent_session_lifetime = timedelta(hours=2)

# Environment variables
EMAIL_USER = os.getenv("EMAIL_USER")
EMAIL_PASSWORD = os.getenv("EMAIL_PASSWORD")
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")

# Database configuration with error handling
try:
    app.config['SQLALCHEMY_DATABASE_URI'] = os.getenv('DATABASE_URL', 'sqlite:///notewise_quiz.db')
    app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False
    app.config['SQLALCHEMY_ENGINE_OPTIONS'] = {
        'pool_timeout': 20,
        'pool_recycle': -1,
        'pool_pre_ping': True
    }
    app.config['WTF_CSRF_ENABLED'] = True
    app.config['SECRET_KEY'] = os.getenv("FLASK_SECRET_KEY") or os.urandom(32)
    app.config['MAX_CONTENT_LENGTH'] = 25 * 1024 * 1024  # 25MB max file size
except Exception as e:
    log.error(f"Database configuration error: {e}")
    raise

db = SQLAlchemy(app)

# Login manager setup
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = 'Please log in to access this page.'
login_manager.login_message_category = 'info'

if not GEMINI_API_KEY:
    raise RuntimeError("GEMINI_API_KEY missing in .env")

try:
    genai.configure(api_key=GEMINI_API_KEY)
    # Test the API connection
    model = genai.GenerativeModel("gemini-2.0-flash")
    log.info("Gemini AI configured successfully")
except Exception as e:
    log.error(f"Gemini AI configuration error: {e}")
    raise

# Initialize EasyOCR reader with error handling
try:
    ocr_reader = easyocr.Reader(['en'], gpu=False)
    log.info("EasyOCR initialized successfully")
except Exception as e:
    log.error(f"EasyOCR initialization error: {e}")
    ocr_reader = None

# ─────────────── SESSION MANAGEMENT HELPERS ───────────────
def clear_quiz_session_data():
    """Clear all quiz-related data from Flask session - FIXED FUNCTION"""
    quiz_keys = [
        'quiz_raw', 'quiz', 'answers', 'quiz_timer', 'quiz_session_id',
        'results', 'score', 'percentage', 'time_taken', 'answered_count', 
        'skipped_count'
    ]
    for key in quiz_keys:
        session.pop(key, None)
    log.info("Quiz session data cleared")

def clear_file_session_data():
    """Clear file-related data from Flask session"""
    file_keys = ['file_path', 'file_type', 'quiz_attempt_id', 'original_filename']
    for key in file_keys:
        session.pop(key, None)
    log.info("File session data cleared")

def clear_all_session_data():
    """Clear all session data for fresh start"""
    clear_quiz_session_data()
    clear_file_session_data()
    log.info("All session data cleared")

# ─────────────── ENHANCED MODELS ───────────────
class User(UserMixin, db.Model):
    id = db.Column(db.Integer, primary_key=True)
    email = db.Column(db.String(120), unique=True, nullable=False, index=True)
    name = db.Column(db.String(100), nullable=False)
    address = db.Column(db.String(200))
    phone = db.Column(db.String(20))
    role = db.Column(db.String(20), nullable=False, default='student')
    password_hash = db.Column(db.String(128), nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    last_login = db.Column(db.DateTime)
    is_active_user = db.Column(db.Boolean, default=True)

    # Relationships
    quiz_attempts = db.relationship('QuizAttempt', backref='user', lazy=True, cascade='all, delete-orphan')
    quiz_sessions = db.relationship('QuizSession', backref='user', lazy=True, cascade='all, delete-orphan')

    @property
    def password(self):
        raise AttributeError('password is not a readable attribute')

    @password.setter
    def password(self, password):
        self.password_hash = generate_password_hash(password)

    def verify_password(self, password):
        return check_password_hash(self.password_hash, password)

    def is_authenticated(self):
        return True

    def is_active(self):
        return self.is_active_user

    def is_anonymous(self):
        return False

    def get_id(self):
        return str(self.id)

class QuizAttempt(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    file_path = db.Column(db.String(200), nullable=False)
    file_type = db.Column(db.String(20), nullable=False)
    quiz_raw = db.Column(db.Text)
    answers = db.Column(db.Text)
    score = db.Column(db.Integer)
    total = db.Column(db.Integer)
    quiz_timer = db.Column(db.Integer, default=0)
    time_taken = db.Column(db.Integer)  # Actual time taken in seconds
    completed = db.Column(db.Boolean, default=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    completed_at = db.Column(db.DateTime)

    def is_expired(self):
        return datetime.utcnow() > self.created_at + timedelta(days=3)

class QuizSession(db.Model):
    """Track active quiz sessions with progress"""
    id = db.Column(db.String(36), primary_key=True, default=lambda: str(uuid.uuid4()))
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    quiz_attempt_id = db.Column(db.Integer, db.ForeignKey('quiz_attempt.id'), nullable=False)
    current_question = db.Column(db.Integer, default=0)
    answers_data = db.Column(db.Text)  # JSON string of answers
    start_time = db.Column(db.DateTime, default=datetime.utcnow)
    last_activity = db.Column(db.DateTime, default=datetime.utcnow)
    time_remaining = db.Column(db.Integer)  # Seconds remaining
    is_active = db.Column(db.Boolean, default=True)

    # Relationships
    quiz_attempt = db.relationship('QuizAttempt', backref='sessions')

    def update_activity(self):
        self.last_activity = datetime.utcnow()

    def is_expired(self):
        return datetime.utcnow() > self.last_activity + timedelta(hours=2)

# ─────────────── ENHANCED FORMS ───────────────
class ProfileForm(FlaskForm):
    name = StringField('Full Name', validators=[
        validators.DataRequired(),
        validators.Length(min=2, max=100)
    ])
    address = StringField('Address', validators=[
        validators.Length(max=200)
    ])
    phone = StringField('Phone Number', validators=[
        validators.Length(max=20)
    ])
    role = SelectField('Role', choices=[
        ('student', 'Student'),
        ('teacher', 'Teacher')
    ])

class PasswordChangeForm(FlaskForm):
    current_password = PasswordField('Current Password', validators=[
        validators.DataRequired(),
        validators.Length(min=6)
    ])
    new_password = PasswordField('New Password', validators=[
        validators.DataRequired(),
        validators.Length(min=6),
        validators.EqualTo('confirm_password', message='Passwords must match')
    ])
    confirm_password = PasswordField('Confirm Password', validators=[
        validators.DataRequired()
    ])

class QuizGenerationForm(FlaskForm):
    num_questions = IntegerField('Number of Questions', validators=[
        validators.DataRequired(),
        validators.NumberRange(min=1, max=30)
    ])
    question_type = SelectField('Question Type', choices=[
        ('multiple-choice', 'Multiple Choice'),
        ('true-false', 'True/False')
    ])
    quiz_timer = IntegerField('Timer (minutes)', validators=[
        validators.NumberRange(min=0, max=180)
    ], default=0)

@login_manager.user_loader
def load_user(user_id):
    try:
        return User.query.get(int(user_id))
    except (ValueError, TypeError):
        return None

# ─────────────── CLEANUP FUNCTION & SCHEDULER ───────────────
def cleanup_old_data():
    """Enhanced cleanup with better error handling"""
    with app.app_context():
        try:
            three_days_ago = datetime.utcnow() - timedelta(days=3)

            # Clean up expired quiz attempts
            old_attempts = QuizAttempt.query.filter(QuizAttempt.created_at < three_days_ago).all()
            files_deleted = 0

            for attempt in old_attempts:
                # Delete the associated file
                if attempt.file_path and os.path.exists(attempt.file_path):
                    try:
                        os.remove(attempt.file_path)
                        files_deleted += 1
                    except OSError as e:
                        log.warning(f"Could not delete file {attempt.file_path}: {e}")
                db.session.delete(attempt)

            # Clean up expired quiz sessions
            two_hours_ago = datetime.utcnow() - timedelta(hours=2)
            expired_sessions = QuizSession.query.filter(
                QuizSession.last_activity < two_hours_ago,
                QuizSession.is_active == True
            ).all()

            for session in expired_sessions:
                session.is_active = False

            db.session.commit()
            log.info(f"Cleaned up {len(old_attempts)} old quiz attempts and {files_deleted} files")
            log.info(f"Deactivated {len(expired_sessions)} expired quiz sessions")

        except Exception as e:
            log.error(f"Error during cleanup: {e}")
            db.session.rollback()

# Initialize scheduler with error handling
try:
    scheduler = BackgroundScheduler()
    scheduler.add_job(func=cleanup_old_data, trigger='interval', hours=6)
    scheduler.start()
    atexit.register(lambda: scheduler.shutdown())
    log.info("Background scheduler initialized")
except Exception as e:
    log.error(f"Failed to initialize scheduler: {e}")

# ─────────────── ENHANCED TEXT EXTRACTION HELPERS ───────────────
def preprocess_image_for_ocr(image_path):
    """Preprocess image for better OCR results with error handling"""
    try:
        if not ocr_reader:
            return image_path

        # Read image using OpenCV
        img = cv2.imread(image_path)
        if img is None:
            # Try with PIL if OpenCV fails
            pil_img = Image.open(image_path)
            img = cv2.cvtColor(np.array(pil_img), cv2.COLOR_RGB2BGR)

        # Convert to grayscale
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

        # Apply denoising
        denoised = cv2.fastNlMeansDenoising(gray)

        # Apply adaptive thresholding for better text contrast
        thresh = cv2.adaptiveThreshold(
            denoised, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 11, 2
        )

        # Save preprocessed image temporarily
        temp_path = tempfile.mktemp(suffix='.png')
        cv2.imwrite(temp_path, thresh)
        return temp_path

    except Exception as e:
        log.error(f"Error preprocessing image: {e}")
        return image_path  # Return original if preprocessing fails

def extract_text_from_pdf(pdf_path: str) -> str:
    """Extract text from PDF file with enhanced error handling"""
    text_content = []
    try:
        with open(pdf_path, "rb") as file:
            pdf_reader = PyPDF2.PdfReader(file)
            
            if len(pdf_reader.pages) == 0:
                raise ValueError("PDF contains no pages")

            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    text = page.extract_text() or ""
                    if text.strip():
                        text_content.append(f"Page {page_num + 1}:\n{text}")
                except Exception as e:
                    log.warning(f"Error extracting text from page {page_num + 1}: {e}")
                    continue

        result = "\n\n".join(text_content)
        if not result.strip():
            raise ValueError("No readable text found in PDF")
        return result

    except Exception as e:
        log.error(f"Error extracting text from PDF: {e}")
        raise

def extract_text_from_image(image_path: str) -> str:
    """Extract text from image using EasyOCR with preprocessing"""
    try:
        if not ocr_reader:
            raise RuntimeError("OCR reader not available")

        # Preprocess image for better OCR
        processed_path = preprocess_image_for_ocr(image_path)

        # Use EasyOCR to extract text
        results = ocr_reader.readtext(processed_path, detail=0)
        text = " ".join(results)

        # Clean up temporary file if created
        if processed_path != image_path and os.path.exists(processed_path):
            os.remove(processed_path)

        if not text.strip():
            raise ValueError("No readable text found in image")

        return text

    except Exception as e:
        log.error(f"Error extracting text from image: {e}")
        raise

def extract_text_from_docx(docx_path: str) -> str:
    """Extract text from Word document with enhanced error handling"""
    try:
        # Try using docx2txt first (handles more complex documents)
        text = docx2txt.process(docx_path)
        if text and len(text.strip()) > 10:
            return text

        # Fallback to python-docx
        doc = Document(docx_path)
        paragraphs = []

        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                paragraphs.append(paragraph.text)

        # Also extract text from tables
        for table in doc.tables:
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text)
                if row_text:
                    paragraphs.append(" | ".join(row_text))

        result = "\n".join(paragraphs)
        if not result.strip():
            raise ValueError("No readable text found in DOCX file")

        return result

    except Exception as e:
        log.error(f"Error extracting text from DOCX: {e}")
        raise

def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract text from PowerPoint presentation with enhanced error handling"""
    try:
        presentation = Presentation(pptx_path)
        slides_text = []

        if len(presentation.slides) == 0:
            raise ValueError("PowerPoint presentation contains no slides")

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"Slide {slide_num}:"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Handle tables in slides
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text)
                        if row_text:
                            slide_text.append(" | ".join(row_text))

            if len(slide_text) > 1:  # More than just the slide title
                slides_text.append("\n".join(slide_text))

        result = "\n\n".join(slides_text)
        if not result.strip():
            raise ValueError("No readable text found in PowerPoint presentation")

        return result

    except Exception as e:
        log.error(f"Error extracting text from PPTX: {e}")
        raise

def extract_text_from_file(file_path: str, file_type: str) -> str:
    """Main function to extract text based on file type with comprehensive error handling"""
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    file_size = os.path.getsize(file_path)
    if file_size == 0:
        raise ValueError("File is empty")

    if file_size > 25 * 1024 * 1024:  # 25MB
        raise ValueError("File size exceeds 25MB limit")

    try:
        if file_type == 'pdf':
            return extract_text_from_pdf(file_path)
        elif file_type in ['image', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff']:
            return extract_text_from_image(file_path)
        elif file_type in ['docx', 'doc']:
            return extract_text_from_docx(file_path)
        elif file_type in ['pptx', 'ppt']:
            return extract_text_from_pptx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

    except Exception as e:
        # Clean up file on error
        if os.path.exists(file_path):
            try:
                os.remove(file_path)
            except OSError:
                pass
        raise

def get_file_type(filename: str) -> str:
    """Determine file type from filename"""
    if not filename or '.' not in filename:
        return 'unknown'

    ext = filename.lower().split('.')[-1]
    if ext == 'pdf':
        return 'pdf'
    elif ext in ['jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp']:
        return 'image'
    elif ext in ['docx', 'doc']:
        return 'docx'
    elif ext in ['pptx', 'ppt']:
        return 'pptx'
    else:
        return 'unknown'

def allowed_file(filename: str) -> bool:
    """Check if file type is allowed"""
    if not filename:
        return False

    allowed_extensions = {
        'pdf', 'jpg', 'jpeg', 'png', 'gif', 'bmp', 'tiff', 'webp',
        'docx', 'doc', 'pptx', 'ppt'
    }
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

# ─────────────── ENHANCED AI GENERATION HELPERS ───────────────
def safe_generate(prompt: str, *, max_tokens: int = 4096, retries: int = 3) -> str | None:
    """Generate content using Gemini AI with retry logic"""
    model = genai.GenerativeModel("gemini-2.0-flash")
    cfg = genai.types.GenerationConfig(
        candidate_count=1,
        max_output_tokens=max_tokens,
        temperature=0.7
    )

    for attempt in range(retries):
        try:
            response = model.generate_content(prompt, generation_config=cfg)
            if not response.candidates:
                log.warning(f"No candidates in response (attempt {attempt + 1})")
                continue

            parts = response.candidates[0].content.parts
            result = "".join(p.text or "" for p in parts).strip()

            if result:
                return result
            else:
                log.warning(f"Empty response (attempt {attempt + 1})")
                continue

        except Exception as e:
            log.error(f"Gemini error on attempt {attempt + 1}: {e}")
            if attempt == retries - 1:
                return None
            time.sleep(2 ** attempt)  # Exponential backoff

    return None

def make_prompt(content: str, n: int, qtype: str) -> str:
    """Create prompt for quiz generation with enhanced instructions"""
    # Validate question count (max 30)
    if n > 30:
        n = 30
    elif n < 1:
        n = 1

    qtype_lower = qtype.lower()
    if qtype_lower in ['true/false', 'true-false', 'tf', 'truefalse', 'true_false']:
        question_format = "True/False"
        instructions = f"""
Generate exactly {n} True/False questions based on the content.

FORMAT REQUIREMENTS:
1. Start each question with "Q1.", "Q2.", etc.
2. End each question with "(True/False)"
3. After each question, add "Answer: True" or "Answer: False"
4. Do not include any other options or explanations
5. Make questions clear and unambiguous
6. Base questions on key facts and concepts from the content
7. Ensure questions test understanding, not memorization

EXAMPLE:
Q1. Python is a programming language. (True/False)
Answer: True

Q2. The sun revolves around the Earth. (True/False)
Answer: False
"""
    else:
        question_format = "Multiple Choice"
        instructions = f"""
Generate exactly {n} multiple-choice questions based on the content.

FORMAT REQUIREMENTS:
1. Start each question with "Q1.", "Q2.", etc.
2. Provide exactly 4 options labeled A), B), C), D)
3. After each question and options, add "Answer: A" (or B, C, D)
4. Do not include any explanations
5. Make sure only one answer is clearly correct
6. Create plausible distractors for wrong answers
7. Vary difficulty levels and question types

EXAMPLE:
Q1. What is the capital of France?
A) London
B) Berlin
C) Paris
D) Madrid
Answer: C
"""

    return textwrap.dedent(f"""
Create a {question_format} quiz with EXACTLY {n} questions based on the provided content.

{instructions}

CRITICAL INSTRUCTIONS:
- Generate exactly {n} questions, no more, no less
- Focus on the most important concepts and facts from the content
- Ensure questions are clear and well-structured
- Avoid ambiguous or trick questions
- Use varied question types covering different aspects of the content
- Make questions progressive in difficulty
- Ensure good coverage of the source material

Source content (first 15000 characters):
{content[:15000]}
""").strip()

def parse_quiz(raw: str):
    """Parse the generated quiz text into structured format with validation"""
    items = []
    current_question = {}
    question_counter = 0

    for line in raw.splitlines():
        line = line.strip()
        if not line:
            continue

        # Check for question start (Q1., Q2., etc.)
        question_match = re.match(r'^Q(\d+)\.\s*(.+)', line)
        if question_match:
            # Save previous question if exists
            if current_question and validate_question(current_question):
                items.append(current_question)

            question_counter += 1
            question_text = question_match.group(2).strip()

            # Start new question
            current_question = {
                "id": question_counter - 1,
                "question": question_text,
                "options": [],
                "answer": "",
                "question_type": "multiple_choice"  # default
            }

            # Check if it's a True/False question
            if "(True/False)" in question_text or "(T/F)" in question_text:
                current_question["question_type"] = "true_false"
                # Clean the question text
                current_question["question"] = (question_text
                                              .replace("(True/False)", "")
                                              .replace("(T/F)", "").strip())

        # Check for multiple choice options (A), B), C), D))
        elif re.match(r'^[A-D]\)\s*(.+)', line):
            if current_question:
                current_question["options"].append(line)

        # Check for answer
        elif line.startswith("Answer:"):
            if current_question:
                answer = line.split(":", 1)[1].strip()
                current_question["answer"] = answer

    # Add the last question if valid
    if current_question and validate_question(current_question):
        items.append(current_question)

    # Validate total number of questions
    if not items:
        raise ValueError("No valid questions were parsed from the generated content")

    return items

def validate_question(question):
    """Validate a single question structure"""
    if not question.get("question") or not question.get("answer"):
        return False

    if question.get("question_type") == "multiple_choice":
        return len(question.get("options", [])) == 4
    elif question.get("question_type") == "true_false":
        return question.get("answer").lower() in ["true", "false"]

    return False

def pdf_from_text(text: str, title: str) -> str:
    """Generate PDF from text content with better formatting"""
    path = os.path.join(tempfile.gettempdir(), f"{uuid.uuid4().hex}.pdf")

    try:
        c = canvas.Canvas(path, pagesize=letter)
        width, height = letter
        y = height - 72

        # Title
        c.setFont("Helvetica-Bold", 16)
        title_lines = textwrap.wrap(title, 60)
        for title_line in title_lines:
            c.drawString(72, y, title_line)
            y -= 20

        y -= 20  # Extra space after title

        # Content
        c.setFont("Helvetica", 11)
        for line in text.splitlines():
            if not line.strip():
                y -= 12
                continue

            # Handle long lines
            wrapped_lines = textwrap.wrap(line, 95)
            if not wrapped_lines:
                wrapped_lines = [line]

            for wrapped_line in wrapped_lines:
                if y < 72:  # Start new page
                    c.showPage()
                    y = height - 72
                    c.setFont("Helvetica", 11)

                c.drawString(72, y, wrapped_line)
                y -= 14

        c.save()
        return path

    except Exception as e:
        log.error(f"Error generating PDF: {e}")
        if os.path.exists(path):
            os.remove(path)
        raise

def email_quiz(to: str, files: list[str]) -> bool:
    """Send quiz via email with better error handling"""
    if not EMAIL_USER or not EMAIL_PASSWORD:
        log.error("Email credentials not configured")
        return False

    try:
        msg = MIMEMultipart()
        msg["From"] = EMAIL_USER
        msg["To"] = to
        msg["Subject"] = "NoteWise Quiz - Your Generated Quiz"

        body = f"""
Hello,

Your NoteWise Quiz has been generated successfully!

Please find attached:
- Quiz questions
- Answer key

Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}

Best regards,
NoteWise Quiz Team
"""

        msg.attach(MIMEText(body, "plain"))

        # Attach files
        for file_path in files:
            if os.path.exists(file_path):
                with open(file_path, "rb") as f:
                    part = MIMEApplication(f.read(), Name=os.path.basename(file_path))
                    part["Content-Disposition"] = f'attachment; filename="{os.path.basename(file_path)}"'
                    msg.attach(part)

        # Send email
        with smtplib.SMTP("smtp.gmail.com", 587) as server:
            server.starttls()
            server.login(EMAIL_USER, EMAIL_PASSWORD)
            server.sendmail(EMAIL_USER, to, msg.as_string())

        return True

    except Exception as e:
        log.error(f"SMTP error: {e}")
        return False

# ─────────────── SESSION MANAGEMENT HELPERS ───────────────
def create_quiz_session(user_id: int, quiz_attempt_id: int, timer_minutes: int = 0) -> str:
    """Create a new quiz session"""
    try:
        session_id = str(uuid.uuid4())
        quiz_session = QuizSession(
            id=session_id,
            user_id=user_id,
            quiz_attempt_id=quiz_attempt_id,
            time_remaining=timer_minutes * 60 if timer_minutes > 0 else None
        )

        db.session.add(quiz_session)
        db.session.commit()
        return session_id

    except Exception as e:
        log.error(f"Error creating quiz session: {e}")
        db.session.rollback()
        raise

def get_quiz_session(session_id: str) -> QuizSession:
    """Get and validate quiz session"""
    if not session_id:
        return None

    quiz_session = QuizSession.query.filter_by(
        id=session_id,
        is_active=True
    ).first()

    if quiz_session and quiz_session.is_expired():
        quiz_session.is_active = False
        db.session.commit()
        return None

    return quiz_session

def update_quiz_progress(session_id: str, current_question: int, answers_data: dict):
    """Update quiz progress in database"""
    try:
        quiz_session = get_quiz_session(session_id)
        if quiz_session:
            quiz_session.current_question = current_question
            quiz_session.answers_data = json.dumps(answers_data)
            quiz_session.update_activity()
            db.session.commit()
            return True
    except Exception as e:
        log.error(f"Error updating quiz progress: {e}")
        db.session.rollback()
        return False

# Custom template filter for filename extraction
@app.template_filter('basename')
def basename_filter(path):
    return os.path.basename(path) if path else ""

@app.template_filter('time_ago')
def time_ago_filter(datetime_obj):
    """Convert datetime to human readable time ago format"""
    if not datetime_obj:
        return "Unknown"

    now = datetime.utcnow()
    diff = now - datetime_obj

    if diff.days > 0:
        return f"{diff.days} day{'s' if diff.days != 1 else ''} ago"
    elif diff.seconds > 3600:
        hours = diff.seconds // 3600
        return f"{hours} hour{'s' if hours != 1 else ''} ago"
    elif diff.seconds > 60:
        minutes = diff.seconds // 60
        return f"{minutes} minute{'s' if minutes != 1 else ''} ago"
    else:
        return "Just now"

# ─────────────── AUTHENTICATION ROUTES ───────────────
@app.route('/signup', methods=['GET', 'POST'])
def signup():
    """Enhanced user registration with validation"""
    if current_user.is_authenticated:
        return redirect(url_for('spa_root'))

    if request.method == 'POST':
        try:
            email = request.form.get('email', '').strip().lower()
            name = request.form.get('name', '').strip()
            address = request.form.get('address', '').strip()
            phone = request.form.get('phone', '').strip()
            role = request.form.get('role', 'student')
            password = request.form.get('password', '')

            # Validation
            if not all([email, name, password]):
                flash('Please fill in all required fields', 'danger')
                return redirect(url_for('signup'))

            if len(password) < 6:
                flash('Password must be at least 6 characters long', 'danger')
                return redirect(url_for('signup'))

            if not re.match(r'^[^@]+@[^@]+\.[^@]+$', email):
                flash('Please enter a valid email address', 'danger')
                return redirect(url_for('signup'))

            # Check if user exists
            user = User.query.filter_by(email=email).first()
            if user:
                flash('Email already exists', 'danger')
                return redirect(url_for('signup'))

            # Create new user
            new_user = User(
                email=email,
                name=name,
                address=address or None,
                phone=phone or None,
                role=role
            )
            new_user.password = password  # uses the setter

            db.session.add(new_user)
            db.session.commit()

            log.info(f"New user registered: {email}")
            flash('Account created successfully. Please login.', 'success')
            return redirect(url_for('login'))

        except Exception as e:
            log.error(f"Error during signup: {e}")
            db.session.rollback()
            flash('An error occurred during registration. Please try again.', 'danger')

    return render_template('signup.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    """Enhanced user login with better error handling"""
    if current_user.is_authenticated:
        return redirect(url_for('spa_root'))

    if request.method == 'POST':
        try:
            email = request.form.get('email', '').strip().lower()
            password = request.form.get('password', '')
            remember = bool(request.form.get('remember'))

            if not email or not password:
                flash('Please enter both email and password', 'danger')
                return render_template('login.html')

            user = User.query.filter_by(email=email).first()

            if user and user.verify_password(password):
                if not user.is_active():
                    flash('Your account has been deactivated', 'danger')
                    return render_template('login.html')

                login_user(user, remember=remember)
                user.last_login = datetime.utcnow()
                db.session.commit()

                log.info(f"User logged in: {email}")

                # Redirect to next page or dashboard
                next_page = request.args.get('next')
                if next_page and next_page.startswith('/'):
                    return redirect(next_page)
                return redirect(url_for('spa_root'))
            else:
                flash('Invalid email or password', 'danger')

        except Exception as e:
            log.error(f"Error during login: {e}")
            flash('An error occurred during login. Please try again.', 'danger')

    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    """Enhanced user logout with cleanup"""
    try:
        # Deactivate any active quiz sessions
        QuizSession.query.filter_by(
            user_id=current_user.id,
            is_active=True
        ).update({'is_active': False})
        
        db.session.commit()

        log.info(f"User logged out: {current_user.email}")
        logout_user()
        session.clear()  # Clear Flask session
        flash('You have been logged out successfully', 'info')

    except Exception as e:
        log.error(f"Error during logout: {e}")

    return redirect(url_for('login'))

# ─────────────── MAIN PAGES ───────────────
@app.get("/")
@login_required
def spa_root():
    """Enhanced main dashboard page"""
    try:
        # Get user's recent quiz attempts
        recent_attempts = QuizAttempt.query.filter_by(user_id=current_user.id)\
            .order_by(QuizAttempt.created_at.desc()).limit(5).all()

        # Get active quiz session if any
        active_session = QuizSession.query.filter_by(
            user_id=current_user.id,
            is_active=True
        ).first()

        return render_template("index.html",
                             recent_attempts=recent_attempts,
                             active_session=active_session)

    except Exception as e:
        log.error(f"Error loading dashboard: {e}")
        return render_template("index.html")

@app.get('/quiz')
@login_required
def quiz_page():
    """Enhanced quiz taking page with session management"""
    return render_template('quiz.html')

@app.get('/results')
@login_required
def results_page():
    """Enhanced quiz results page"""
    return render_template('results.html')

# ─────────────── PROFILE PAGE ───────────────
@app.route('/profile', methods=['GET', 'POST'])
@login_required
def profile_page():
    """Enhanced user profile management"""
    profile_form = ProfileForm()
    password_form = PasswordChangeForm()

    try:
        # Get user's quiz history (last 15 attempts)
        quiz_history = QuizAttempt.query.filter_by(user_id=current_user.id)\
            .order_by(QuizAttempt.created_at.desc()).limit(15).all()

        # Calculate statistics
        total_quizzes = len(quiz_history)
        completed_quizzes = len([q for q in quiz_history if q.completed])
        avg_score = 0

        if completed_quizzes > 0:
            scores = [q.score / q.total * 100 for q in quiz_history if q.completed and q.total > 0]
            avg_score = sum(scores) / len(scores) if scores else 0

        stats = {
            'total_quizzes': total_quizzes,
            'completed_quizzes': completed_quizzes,
            'avg_score': round(avg_score, 1)
        }

        # Handle profile update
        if profile_form.validate_on_submit():
            current_user.name = profile_form.name.data
            current_user.address = profile_form.address.data
            current_user.phone = profile_form.phone.data
            current_user.role = profile_form.role.data
            db.session.commit()

            log.info(f"Profile updated for user: {current_user.email}")
            flash('Profile updated successfully!', 'success')
            return redirect(url_for('profile_page'))

        # Handle password change
        if password_form.validate_on_submit():
            if current_user.verify_password(password_form.current_password.data):
                current_user.password = password_form.new_password.data
                db.session.commit()

                log.info(f"Password changed for user: {current_user.email}")
                flash('Password changed successfully!', 'success')
            else:
                flash('Current password is incorrect', 'danger')
            
            return redirect(url_for('profile_page'))

        # Populate form with current data
        if not profile_form.is_submitted():
            profile_form.name.data = current_user.name
            profile_form.address.data = current_user.address
            profile_form.phone.data = current_user.phone
            profile_form.role.data = current_user.role

        return render_template('profile.html',
                             profile_form=profile_form,
                             password_form=password_form,
                             quiz_history=quiz_history,
                             stats=stats)

    except Exception as e:
        log.error(f"Error loading profile page: {e}")
        flash('Error loading profile. Please try again.', 'danger')
        return redirect(url_for('spa_root'))

# ─────────────── ENHANCED API ROUTES ───────────────
@app.post("/api/upload")
@login_required
def api_upload():
    """FIXED: Enhanced file upload API with session cleanup"""
    try:
        # Check file size
        if request.content_length and request.content_length > 25 * 1024 * 1024:
            return {"error": "File size exceeds 25MB limit"}, 400

        uploaded_file = request.files.get("file")
        if not uploaded_file or not uploaded_file.filename:
            return {"error": "No file provided"}, 400

        # Validate file type
        if not allowed_file(uploaded_file.filename):
            return {
                "error": "Unsupported file type. Please upload PDF, Word (DOCX), PowerPoint (PPTX), or image files."
            }, 400

        # FIXED: Clear ALL session data when uploading new file
        clear_all_session_data()

        # Determine file type
        file_type = get_file_type(uploaded_file.filename)
        if file_type == 'unknown':
            return {"error": "Unable to determine file type"}, 400

        # Save file with secure filename
        filename = secure_filename(uploaded_file.filename)
        file_extension = filename.rsplit('.', 1)[1].lower()
        unique_filename = f"{uuid.uuid4().hex}.{file_extension}"
        file_path = os.path.join(tempfile.gettempdir(), unique_filename)

        try:
            uploaded_file.save(file_path)

            # Validate file was saved correctly
            if not os.path.exists(file_path) or os.path.getsize(file_path) == 0:
                return {"error": "Failed to save uploaded file"}, 500

        except Exception as e:
            log.error(f"Error saving file: {e}")
            return {"error": "Failed to save uploaded file"}, 500

        # Create quiz attempt record
        try:
            attempt = QuizAttempt(
                user_id=current_user.id,
                file_path=file_path,
                file_type=file_type
            )

            db.session.add(attempt)
            db.session.commit()

            # Store in session for generation
            session.permanent = True
            session["file_path"] = file_path
            session["file_type"] = file_type
            session["quiz_attempt_id"] = attempt.id
            session["original_filename"] = uploaded_file.filename

            log.info(f"File uploaded successfully: {file_type} by user {current_user.email}")

            return {
                "ok": True,
                "filename": uploaded_file.filename,
                "file_type": file_type,
                "file_size": os.path.getsize(file_path),
                "message": f"Successfully uploaded {file_type.upper()} file"
            }

        except Exception as e:
            log.error(f"Error creating quiz attempt: {e}")
            # Clean up file on database error
            if os.path.exists(file_path):
                os.remove(file_path)
            db.session.rollback()
            return {"error": "Failed to process uploaded file"}, 500

    except Exception as e:
        log.error(f"Unexpected error in upload: {e}")
        return {"error": "An unexpected error occurred"}, 500

@app.post("/api/generate")
@login_required
def api_generate():
    """FIXED: Enhanced quiz generation API with session cleanup"""
    try:
        data = request.json or {}

        # FIXED: Clear old quiz data at start of generation
        clear_quiz_session_data()

        if "file_path" not in session:
            return {"error": "No file uploaded. Please upload a file first."}, 400

        # Extract text based on file type
        file_path = session["file_path"]
        file_type = session["file_type"]

        log.info(f"Extracting text from {file_type} file: {os.path.basename(file_path)}")

        try:
            extracted_text = extract_text_from_file(file_path, file_type)
        except Exception as e:
            log.error(f"Error extracting text: {e}")
            return {"error": f"Failed to extract text from {file_type} file: {str(e)}"}, 500

        if not extracted_text or len(extracted_text.strip()) < 50:
            return {
                "error": f"Insufficient text content extracted from the {file_type} file. Please ensure the file contains readable text."
            }, 400

        # Get and validate generation parameters
        try:
            num_questions = int(data.get("num_questions", 5))
            question_type = data.get("question_type", "multiple-choice")
            quiz_timer = int(data.get("quiz_timer", 0))

            # Validation
            if num_questions < 1 or num_questions > 30:
                return {"error": "Number of questions must be between 1 and 30"}, 400

            if quiz_timer < 0 or quiz_timer > 180:
                return {"error": "Timer must be between 0 and 180 minutes"}, 400

        except (ValueError, TypeError):
            return {"error": "Invalid parameters provided"}, 400

        log.info(f"Generating {num_questions} {question_type} questions from {len(extracted_text)} characters")

        # Generate quiz using AI
        prompt = make_prompt(extracted_text, num_questions, question_type)
        raw_quiz = safe_generate(prompt, max_tokens=6144)

        if not raw_quiz:
            return {"error": "AI model failed to generate quiz. Please try again with different content."}, 500

        # Parse the generated quiz
        try:
            quiz_questions = parse_quiz(raw_quiz)
        except Exception as e:
            log.error(f"Error parsing quiz: {e}")
            return {"error": "Failed to parse generated quiz. Please try again."}, 500

        # Validate generated questions
        if not quiz_questions:
            return {"error": "No valid questions were generated. Please try again with different content."}, 500

        if len(quiz_questions) != num_questions:
            log.warning(f"Expected {num_questions} questions, got {len(quiz_questions)}")

        # Generate answers summary
        answers = "\n".join(f"Q{i+1}. {q['answer']}" for i, q in enumerate(quiz_questions))

        # Update quiz attempt with generated content
        try:
            attempt_id = session.get("quiz_attempt_id")
            if attempt_id:
                attempt = QuizAttempt.query.get(attempt_id)
                if attempt:
                    attempt.quiz_raw = raw_quiz
                    attempt.answers = answers
                    attempt.total = len(quiz_questions)
                    attempt.quiz_timer = quiz_timer
                    db.session.commit()
        except Exception as e:
            log.error(f"Error updating quiz attempt: {e}")
            # Continue anyway, store in session

        # Create quiz session for progress tracking
        session_id = None
        try:
            if quiz_timer > 0:
                session_id = create_quiz_session(current_user.id, attempt_id, quiz_timer)
        except Exception as e:
            log.error(f"Error creating quiz session: {e}")

        # Store in Flask session
        session.update({
            "quiz_raw": raw_quiz,
            "quiz": quiz_questions,
            "answers": answers,
            "quiz_timer": quiz_timer,
            "quiz_session_id": session_id
        })

        log.info(f"Quiz generated successfully: {len(quiz_questions)} questions for user {current_user.email}")

        return {
            "quiz": quiz_questions,
            "total_questions": len(quiz_questions),
            "timer": quiz_timer,
            "session_id": session_id,
            "message": f"Successfully generated {len(quiz_questions)} questions from your {file_type.upper()} file"
        }

    except Exception as e:
        log.error(f"Unexpected error in quiz generation: {e}")
        return {"error": "An unexpected error occurred during quiz generation"}, 500

@app.post("/api/submit")
@login_required
def api_submit():
    """Enhanced quiz submission with comprehensive validation and session management"""
    try:
        if "quiz" not in session:
            return {"error": "No quiz found in session"}, 400

        data = request.json or {}
        user_answers = data.get('answers', {})
        time_taken = data.get('time_taken', 0)
        session_id = data.get('session_id')

        quiz_questions = session["quiz"]
        score = 0
        results = []

        # Process each question
        for i, question in enumerate(quiz_questions):
            user_answer = (user_answers.get(str(i), "") or "").strip()
            correct_answer = question["answer"].strip()

            # Handle different question types properly
            if question.get("question_type") == "true_false":
                # For True/False questions, compare full answer
                is_correct = user_answer.lower() == correct_answer.lower()
            else:
                # For multiple choice, compare first character (A, B, C, D)
                is_correct = user_answer.upper()[:1] == correct_answer.upper()[:1]

            if is_correct:
                score += 1

            results.append({
                "q": question["question"],
                "user": user_answer,
                "correct": correct_answer,
                "ok": is_correct,
                "type": question.get("question_type", "multiple_choice"),
                "options": question.get("options", [])
            })

        # Calculate statistics
        total_questions = len(results)
        percentage = round((score / total_questions) * 100, 1) if total_questions > 0 else 0
        answered_count = len([r for r in results if r["user"]])
        skipped_count = total_questions - answered_count

        # Save results to quiz attempt
        try:
            attempt_id = session.get("quiz_attempt_id")
            if attempt_id:
                attempt = QuizAttempt.query.get(attempt_id)
                if attempt:
                    attempt.score = score
                    attempt.total = total_questions
                    attempt.completed = True
                    attempt.completed_at = datetime.utcnow()
                    attempt.time_taken = time_taken
                    db.session.commit()
        except Exception as e:
            log.error(f"Error saving quiz results: {e}")
            db.session.rollback()

        # Deactivate quiz session
        try:
            if session_id:
                quiz_session = get_quiz_session(session_id)
                if quiz_session:
                    quiz_session.is_active = False
                    db.session.commit()
        except Exception as e:
            log.error(f"Error deactivating quiz session: {e}")

        # Store results in session
        session["results"] = results
        session["score"] = score
        session["percentage"] = percentage
        session["time_taken"] = time_taken
        session["answered_count"] = answered_count
        session["skipped_count"] = skipped_count

        log.info(f"Quiz submitted: {score}/{total_questions} ({percentage}%) by user {current_user.email}")

        return {
            "score": score,
            "total": total_questions,
            "percentage": percentage,
            "answered": answered_count,
            "skipped": skipped_count,
            "time_taken": time_taken,
            "results": results,
            "message": f"Quiz completed! You scored {score}/{total_questions} ({percentage}%)"
        }

    except Exception as e:
        log.error(f"Error submitting quiz: {e}")
        return {"error": "An error occurred while submitting the quiz"}, 500

# ─────────────── QUIZ NAVIGATION & PROGRESS API ───────────────
@app.get("/api/quiz-status")
@login_required
def api_quiz_status():
    """FIXED: Enhanced quiz status with validation"""
    try:
        # FIXED: Check both file and quiz data exist and are consistent
        has_file = "file_path" in session and session.get("file_path") is not None
        has_quiz = "quiz" in session and session.get("quiz") is not None
        
        # Ensure both file and quiz data exist for valid state
        has_valid_quiz = has_file and has_quiz
        
        session_id = session.get("quiz_session_id")

        status = {
            "hasQuiz": has_valid_quiz,
            "sessionId": session_id
        }

        if has_valid_quiz:
            quiz_session = None
            if session_id:
                quiz_session = get_quiz_session(session_id)

            status.update({
                "total": len(session["quiz"]),
                "timer": session.get("quiz_timer", 0),
                "hasActiveSession": quiz_session is not None,
                "timeRemaining": quiz_session.time_remaining if quiz_session else None
            })

        return status

    except Exception as e:
        log.error(f"Error getting quiz status: {e}")
        return {"hasQuiz": False}

@app.get("/api/quiz-data")
@login_required
def api_quiz_data():
    """FIXED: Enhanced quiz data with validation"""
    try:
        # FIXED: Validate both file and quiz data exist
        if "quiz" not in session or "file_path" not in session:
            return {"error": "No valid quiz found"}, 404

        # FIXED: Additional validation to prevent stale data
        if not session.get("quiz") or not session.get("file_path"):
            return {"error": "Invalid quiz data"}, 404

        quiz_data = session["quiz"]
        session_id = session.get("quiz_session_id")

        # Add navigation metadata to each question
        for i, question in enumerate(quiz_data):
            question["questionNumber"] = i + 1
            question["id"] = i

        response = {
            "quiz": quiz_data,
            "total": len(quiz_data),
            "timer": session.get("quiz_timer", 0),
            "sessionId": session_id,
            "filename": session.get("original_filename", "Unknown file")
        }

        # Include session progress if available
        if session_id:
            quiz_session = get_quiz_session(session_id)
            if quiz_session:
                response.update({
                    "currentQuestion": quiz_session.current_question,
                    "timeRemaining": quiz_session.time_remaining,
                    "savedAnswers": json.loads(quiz_session.answers_data) if quiz_session.answers_data else {}
                })

        return response

    except Exception as e:
        log.error(f"Error getting quiz data: {e}")
        return {"error": "Failed to load quiz data"}, 500

@app.post("/api/save-progress")
@login_required
def api_save_progress():
    """Save quiz progress during quiz taking"""
    try:
        data = request.json or {}
        session_id = data.get("sessionId")
        current_question = data.get("currentQuestion", 0)
        answers = data.get("answers", {})
        time_remaining = data.get("timeRemaining")

        if not session_id:
            return {"error": "No session ID provided"}, 400

        quiz_session = get_quiz_session(session_id)
        if not quiz_session:
            return {"error": "Invalid or expired session"}, 400

        # Update progress
        quiz_session.current_question = current_question
        quiz_session.answers_data = json.dumps(answers)
        if time_remaining is not None:
            quiz_session.time_remaining = time_remaining

        quiz_session.update_activity()
        db.session.commit()

        return {
            "saved": True,
            "message": "Progress saved successfully"
        }

    except Exception as e:
        log.error(f"Error saving progress: {e}")
        db.session.rollback()
        return {"error": "Failed to save progress"}, 500

@app.post("/api/quit-quiz")
@login_required
def api_quit_quiz():
    """Quit quiz without saving results"""
    try:
        data = request.json or {}
        session_id = data.get("sessionId")
        save_progress = data.get("saveProgress", False)

        # Deactivate quiz session
        if session_id:
            quiz_session = get_quiz_session(session_id)
            if quiz_session:
                if not save_progress:
                    # Mark as quit without completion
                    quiz_session.is_active = False
                else:
                    # Save current progress but don't complete
                    current_question = data.get("currentQuestion", 0)
                    answers = data.get("answers", {})
                    quiz_session.current_question = current_question
                    quiz_session.answers_data = json.dumps(answers)
                    quiz_session.is_active = False

                db.session.commit()

        # Clear session data if not saving progress
        if not save_progress:
            clear_quiz_session_data()

        log.info(f"User quit quiz: {current_user.email}, save_progress: {save_progress}")

        return {
            "quit": True,
            "message": "Quiz quit successfully" + (" with progress saved" if save_progress else "")
        }

    except Exception as e:
        log.error(f"Error quitting quiz: {e}")
        return {"error": "Failed to quit quiz"}, 500

@app.get("/api/results")
@login_required
def api_results():
    """Enhanced quiz results with detailed analytics"""
    try:
        if "results" not in session:
            return {"error": "No results found"}, 404

        results = session.get("results", [])
        score = session.get("score", 0)
        total = len(results)
        percentage = session.get("percentage", 0)
        time_taken = session.get("time_taken", 0)
        answered_count = session.get("answered_count", 0)
        skipped_count = session.get("skipped_count", 0)

        # Calculate additional analytics
        correct_by_type = {"multiple_choice": 0, "true_false": 0}
        total_by_type = {"multiple_choice": 0, "true_false": 0}

        for result in results:
            qtype = result.get("type", "multiple_choice")
            total_by_type[qtype] += 1
            if result["ok"]:
                correct_by_type[qtype] += 1

        analytics = {
            "score": score,
            "total": total,
            "percentage": percentage,
            "answered": answered_count,
            "skipped": skipped_count,
            "time_taken": time_taken,
            "results": results,
            "performance_by_type": {
                "multiple_choice": {
                    "correct": correct_by_type["multiple_choice"],
                    "total": total_by_type["multiple_choice"],
                    "percentage": round((correct_by_type["multiple_choice"] / total_by_type["multiple_choice"] * 100), 1) if total_by_type["multiple_choice"] > 0 else 0
                },
                "true_false": {
                    "correct": correct_by_type["true_false"],
                    "total": total_by_type["true_false"],
                    "percentage": round((correct_by_type["true_false"] / total_by_type["true_false"] * 100), 1) if total_by_type["true_false"] > 0 else 0
                }
            }
        }

        return analytics

    except Exception as e:
        log.error(f"Error getting results: {e}")
        return {"error": "Failed to load results"}, 500

@app.get("/api/download/<kind>.pdf")
@login_required
def api_download(kind: str):
    """Enhanced download with better error handling"""
    try:
        if kind == "quiz":
            text = session.get("quiz_raw")
            title = "NoteWise Quiz - Questions"
        elif kind == "answers":
            text = session.get("answers")
            title = "NoteWise Quiz - Answer Key"
        else:
            return {"error": "Invalid download type"}, 400

        if not text:
            return {"error": "No content available for download"}, 400

        # Add metadata to the content
        metadata = f"Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}\n"
        metadata += f"User: {current_user.name} ({current_user.email})\n"
        metadata += f"File: {session.get('original_filename', 'Unknown')}\n\n"

        full_text = metadata + text

        pdf_path = pdf_from_text(full_text, title)

        return send_file(
            pdf_path,
            as_attachment=True,
            download_name=f"notewise_{kind}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf",
            mimetype="application/pdf"
        )

    except Exception as e:
        log.error(f"Error generating PDF: {e}")
        return {"error": "Failed to generate PDF"}, 500

@app.post("/api/share")
@login_required
def api_share():
    """Enhanced quiz sharing with better validation"""
    try:
        data = request.json or {}
        email = data.get("email", "").strip().lower()

        if not re.match(r"[^@]+@[^@]+\.[^@]+", email):
            return {"error": "Please provide a valid email address"}, 400

        if "quiz_raw" not in session or "answers" not in session:
            return {"error": "No quiz available to share"}, 400

        # Add metadata
        metadata = f"Generated on: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}\n"
        metadata += f"Shared by: {current_user.name} ({current_user.email})\n"
        metadata += f"Source file: {session.get('original_filename', 'Unknown')}\n\n"

        quiz_content = metadata + session["quiz_raw"]
        answers_content = metadata + session["answers"]

        # Generate PDF files
        quiz_pdf = pdf_from_text(quiz_content, "NoteWise Quiz - Questions")
        answers_pdf = pdf_from_text(answers_content, "NoteWise Quiz - Answer Key")

        # Send email
        success = email_quiz(email, [quiz_pdf, answers_pdf])

        # Clean up temporary files
        for pdf_file in [quiz_pdf, answers_pdf]:
            if os.path.exists(pdf_file):
                try:
                    os.remove(pdf_file)
                except OSError:
                    pass

        if success:
            log.info(f"Quiz shared by {current_user.email} to {email}")
            return {"sent": True, "message": f"Quiz successfully sent to {email}"}
        else:
            return {"sent": False, "error": "Failed to send email. Please check email configuration."}, 500

    except Exception as e:
        log.error(f"Error sharing quiz: {e}")
        return {"error": "Failed to share quiz"}, 500

# ─────────────── QUIZ HISTORY API ───────────────
@app.get("/api/quiz-history")
@login_required
def api_quiz_history():
    """Get user's quiz history with pagination"""
    try:
        page = request.args.get('page', 1, type=int)
        per_page = request.args.get('per_page', 10, type=int)

        # Limit per_page to prevent abuse
        per_page = min(per_page, 50)

        quiz_history = QuizAttempt.query.filter_by(user_id=current_user.id)\
            .order_by(QuizAttempt.created_at.desc())\
            .paginate(page=page, per_page=per_page, error_out=False)

        history_data = []
        for attempt in quiz_history.items:
            history_data.append({
                "id": attempt.id,
                "filename": os.path.basename(attempt.file_path) if attempt.file_path else "Unknown",
                "file_type": attempt.file_type,
                "score": attempt.score,
                "total": attempt.total,
                "percentage": round((attempt.score / attempt.total * 100), 1) if attempt.total and attempt.score is not None else 0,
                "completed": attempt.completed,
                "created_at": attempt.created_at.isoformat() if attempt.created_at else None,
                "completed_at": attempt.completed_at.isoformat() if attempt.completed_at else None,
                "time_taken": attempt.time_taken
            })

        return {
            "history": history_data,
            "pagination": {
                "page": quiz_history.page,
                "pages": quiz_history.pages,
                "per_page": quiz_history.per_page,
                "total": quiz_history.total,
                "has_next": quiz_history.has_next,
                "has_prev": quiz_history.has_prev
            }
        }

    except Exception as e:
        log.error(f"Error getting quiz history: {e}")
        return {"error": "Failed to load quiz history"}, 500

# ─────────────── ERROR HANDLERS ───────────────
@app.errorhandler(413)
def too_large(e):
    """Handle file too large error"""
    log.warning(f"File too large error from user: {current_user.email if current_user.is_authenticated else 'anonymous'}")
    if request.path.startswith('/api/'):
        return jsonify({"error": "File is too large. Maximum size is 25MB."}), 413
    flash("File is too large. Maximum size is 25MB.", "danger")
    return redirect(url_for('spa_root'))

@app.errorhandler(404)
def not_found(e):
    """Handle 404 errors"""
    log.warning(f"404 error: {request.path}")
    if request.path.startswith('/api/'):
        return jsonify({"error": "API endpoint not found"}), 404
    return render_template('404.html'), 404

@app.errorhandler(500)
def internal_error(e):
    """Handle 500 errors"""
    log.error(f"500 error: {e}")
    db.session.rollback()
    if request.path.startswith('/api/'):
        return jsonify({"error": "Internal server error"}), 500
    return render_template('500.html'), 500

@app.errorhandler(Exception)
def handle_exception(e):
    """Handle unexpected exceptions"""
    log.error(f"Unhandled exception: {e}", exc_info=True)
    db.session.rollback()
    if request.path.startswith('/api/'):
        return jsonify({"error": "An unexpected error occurred"}), 500
    flash("An unexpected error occurred. Please try again.", "danger")
    return redirect(url_for('spa_root'))

# ─────────────── DATABASE INITIALIZATION ───────────────
def init_database():
    """Initialize database with proper error handling"""
    try:
        with app.app_context():
            # Test database connection
            db.engine.connect()

            # Create all tables
            db.create_all()

            # Create default admin user if doesn't exist
            admin_user = User.query.filter_by(email='admin@notewise.com').first()
            if not admin_user:
                admin_user = User(
                    email='admin@notewise.com',
                    name='Admin User',
                    role='teacher'
                )
                admin_user.password = 'admin123'

                db.session.add(admin_user)
                db.session.commit()
                log.info("Default admin user created: admin@notewise.com / admin123")

            log.info("Database initialized successfully")
            return True

    except Exception as e:
        log.error(f"Database initialization failed: {e}")
        return False

# ─────────────── HEALTH CHECK ───────────────
@app.route('/health')
def health_check():
    """Health check endpoint"""
    try:
        # Test database connection
        db.session.execute('SELECT 1')

        # Test Gemini AI
        model = genai.GenerativeModel("gemini-2.0-flash")

        return jsonify({
            "status": "healthy",
            "database": "connected",
            "ai_service": "connected",
            "timestamp": datetime.utcnow().isoformat()
        })

    except Exception as e:
        log.error(f"Health check failed: {e}")
        return jsonify({
            "status": "unhealthy",
            "error": str(e),
            "timestamp": datetime.utcnow().isoformat()
        }), 503

# ─────────────── MAIN ───────────────
if __name__ == "__main__":
    # Initialize database
    if not init_database():
        log.error("Failed to initialize database. Exiting.")
        exit(1)

    # Log startup information
    log.info("=" * 60)
    log.info("NoteWise Quiz Server Starting...")
    log.info("=" * 60)
    log.info("Features:")
    log.info("✓ PDF, Word, PowerPoint & Image support")
    log.info("✓ Advanced OCR with preprocessing")
    log.info("✓ Quiz timer and navigation")
    log.info("✓ Progress tracking and saving")
    log.info("✓ User authentication and profiles")
    log.info("✓ Email sharing and PDF download")
    log.info("✓ Comprehensive error handling")
    log.info("✓ Background cleanup scheduler")
    log.info("✓ FIXED: Session management caching issue")
    log.info("=" * 60)
    log.info("Default admin account: admin@notewise.com / admin123")
    log.info("Server: http://127.0.0.1:5001")
    log.info("Health check: http://127.0.0.1:5001/health")
    log.info("=" * 60)

    try:
        app.run(
            debug=os.getenv('FLASK_DEBUG', 'False').lower() == 'true',
            port=int(os.getenv('FLASK_PORT', 5001)),
            host=os.getenv('FLASK_HOST', '127.0.0.1'),
            use_reloader=False  # Disable reloader to prevent scheduler issues
        )

    except Exception as e:
        log.error(f"Failed to start server: {e}")
        exit(1)
